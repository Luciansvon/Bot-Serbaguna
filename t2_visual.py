import os
import httpx
import base64
import json
import csv
import io
from pathlib import Path
from crewai import Agent
from crewai.tools import BaseTool
from crewai_tools import FileReadTool
from config import visual_llm, OUTPUT_DIR

# ============================================================
# Helper: Gemini Vision via OpenRouter (dipakai OCR PDF scan)
# ============================================================
def _gemini_vision(image_b64: str, content_type: str, prompt: str, max_tokens: int = 1500) -> str:
    from openai import OpenAI
    client = OpenAI(
        api_key=os.environ.get("OPENROUTER_API_KEY"),
        base_url="https://openrouter.ai/api/v1"
    )
    result = client.chat.completions.create(
        model="google/gemini-3-flash-preview",
        messages=[{
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": f"data:{content_type};base64,{image_b64}"}},
                {"type": "text", "text": prompt},
            ],
        }],
        max_tokens=max_tokens,
    )
    return result.choices[0].message.content

# ============================================================
# Tool 1: Download File dari URL (Discord Attachment)
# ============================================================
class FileDownloader(BaseTool):
    name: str = "File Downloader Tool"
    description: str = """Download file dari URL (Discord attachment, link drive, dll).
    Simpan ke folder outputs/ dan kembalikan path lokal.
    Input: URL file yang ingin didownload."""

    def _run(self, url: str) -> str:
        try:
            resp = httpx.get(url, timeout=30, follow_redirects=True)
            # Deteksi ekstensi dari Content-Type atau URL
            content_type = resp.headers.get("content-type", "")
            ext = self._get_extension(url, content_type)
            
            filename = f"download_{os.urandom(4).hex()}{ext}"
            filepath = OUTPUT_DIR / filename
            filepath.write_bytes(resp.content)
            
            return f"SUCCESS|{filepath}|File didownload: {filename} ({len(resp.content)} bytes)"
        except Exception as e:
            return f"FAILED|Gagal download: {e}"

    def _get_extension(self, url: str, content_type: str) -> str:
        mapping = {
            "application/pdf": ".pdf",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
            "text/plain": ".txt",
            "text/csv": ".csv",
            "text/markdown": ".md",
            "image/png": ".png",
            "image/jpeg": ".jpg",
            "image/webp": ".webp",
            "audio/mpeg": ".mp3",
            "audio/mp3": ".mp3",
            "audio/wav": ".wav",
            "audio/x-wav": ".wav",
            "audio/ogg": ".ogg",
            "audio/opus": ".opus",
            "audio/mp4": ".m4a",
            "audio/flac": ".flac",
        }
        for ct, ext in mapping.items():
            if ct in content_type:
                return ext
        # Fallback dari URL
        for ext in [".pdf", ".xlsx", ".docx", ".pptx", ".txt", ".csv", ".md",
                    ".png", ".jpg", ".jpeg", ".webp",
                    ".mp3", ".wav", ".m4a", ".ogg", ".opus", ".flac"]:
            if ext in url.lower():
                return ext
        return ".bin"

# ============================================================
# Tool 2: Baca TXT / CSV / MD
# ============================================================
class TextFileReader(BaseTool):
    name: str = "Text File Reader Tool"
    description: str = """Baca dan ekstrak konten dari file teks (TXT, CSV, MD).
    Input: path lokal file yang ingin dibaca."""

    def _run(self, filepath: str) -> str:
        try:
            p = Path(filepath)
            if not p.exists():
                return f"FAILED|File tidak ditemukan: {filepath}"
            
            if p.suffix.lower() == ".csv":
                return self._read_csv(p)
            
            # TXT, MD, dan lainnya
            content = p.read_text(encoding="utf-8", errors="ignore")
            preview = content[:15000]
            if len(content) > 15000:
                preview += f"\n\n... (terpotong di karakter ke-15000, total {len(content)} karakter)"
            return f"=== Isi File: {p.name} ===\n\n{preview}"
        except Exception as e:
            return f"FAILED|Gagal baca file: {e}"

    def _read_csv(self, p: Path) -> str:
        try:
            with open(p, "r", encoding="utf-8") as f:
                reader = csv.reader(f)
                rows = list(reader)
            
            output = [f"=== CSV: {p.name} ===\n"]
            output.append(f"Total baris: {len(rows)} | Total kolom: {len(rows[0]) if rows else 0}\n")

            # Header + 50 baris pertama
            for i, row in enumerate(rows[:51]):
                prefix = "HEADER: " if i == 0 else f"Row {i}: "
                output.append(f"{prefix}{' | '.join(row[:15])}")

            if len(rows) > 51:
                output.append(f"\n... (terpotong, {len(rows)-51} baris lagi)")
            
            return "\n".join(output)
        except Exception as e:
            return f"FAILED|Gagal baca CSV: {e}"

# ============================================================
# Tool 3: Baca PDF (dengan Tabel Support)
# ============================================================
class PDFReader(BaseTool):
    name: str = "PDF Reader Tool"
    description: str = """Baca dan ekstrak teks + tabel dari file PDF.
    Support: katalog, manual, laporan, invoice, dll.
    Input: path lokal file PDF."""

    def _run(self, filepath: str) -> str:
        try:
            import pdfplumber

            p = Path(filepath)
            if not p.exists():
                return f"FAILED|File tidak ditemukan: {filepath}"

            output = [f"=== PDF: {p.name} ===\n"]
            MAX_PAGES = 20

            with pdfplumber.open(p) as pdf:
                output.append(f"Total halaman: {len(pdf.pages)}\n")

                # Deteksi PDF scan: kalau total teks layer < 100 char dari 3 halaman pertama, OCR mode aktif
                sample = "".join((pg.extract_text() or "") for pg in pdf.pages[:3]).strip()
                is_scanned = len(sample) < 100
                if is_scanned:
                    output.append("⚠️ PDF terdeteksi sebagai SCAN (no text layer). Aktifkan OCR via Gemini Vision...\n")

                for i, page in enumerate(pdf.pages[:MAX_PAGES]):
                    output.append(f"--- Halaman {i+1} ---")

                    if is_scanned:
                        try:
                            page_img = page.to_image(resolution=150).original
                            buf = io.BytesIO()
                            page_img.save(buf, format='PNG')
                            b64 = base64.b64encode(buf.getvalue()).decode('utf-8')
                            ocr_text = _gemini_vision(
                                b64, "image/png",
                                "OCR halaman PDF ini. Ekstrak SEMUA teks persis seperti aslinya, termasuk tabel. Pertahankan struktur baris. Output: hanya teks, tanpa komentar atau markdown.",
                                max_tokens=2500
                            )
                            output.append(f"TEKS (OCR):\n{ocr_text}")
                        except Exception as e:
                            output.append(f"OCR gagal di halaman {i+1}: {e}")
                    else:
                        text = page.extract_text()
                        if text:
                            output.append(f"TEKS:\n{text[:1500]}")

                        tables = page.extract_tables()
                        if tables:
                            output.append(f"\nTABEL ({len(tables)} tabel):")
                            for j, table in enumerate(tables[:2]):
                                output.append(f"  Tabel {j+1}:")
                                for row in table[:5]:
                                    output.append(f"    {' | '.join(str(c) for c in row)}")

                    output.append("")

                if len(pdf.pages) > MAX_PAGES:
                    output.append(f"... (terpotong, {len(pdf.pages)-MAX_PAGES} halaman lagi)")

            return "\n".join(output)
        except ImportError:
            return "FAILED|pdfplumber belum terinstall. Jalankan: pip install pdfplumber"
        except Exception as e:
            return f"FAILED|Gagal baca PDF: {e}"

# ============================================================
# Tool 4: Baca Excel (XLSX)
# ============================================================
class ExcelReader(BaseTool):
    name: str = "Excel Reader Tool"
    description: str = """Baca dan ekstrak data dari file Excel (.xlsx).
    Support: multiple sheets, formula, tabel, dan chart data.
    Input: path lokal file Excel."""

    def _run(self, filepath: str) -> str:
        try:
            import openpyxl
            
            p = Path(filepath)
            if not p.exists():
                return f"FAILED|File tidak ditemukan: {filepath}"
            
            wb = openpyxl.load_workbook(p, data_only=True)
            output = [f"=== Excel: {p.name} ===\n"]
            output.append(f"Total sheets: {len(wb.sheetnames)} → {wb.sheetnames}\n")

            MAX_SHEETS = 10
            MAX_ROWS = 50
            MAX_COLS = 15

            for sheet_name in wb.sheetnames[:MAX_SHEETS]:
                ws = wb[sheet_name]
                output.append(f"--- Sheet: {sheet_name} ---")
                output.append(f"Dimensi: {ws.dimensions} | Max row: {ws.max_row} | Max col: {ws.max_column}\n")

                rows = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= MAX_ROWS:
                        break
                    row_data = [str(cell) if cell is not None else "" for cell in row[:MAX_COLS]]
                    rows.append(" | ".join(row_data))

                output.append("\n".join(rows))

                if ws.max_row > MAX_ROWS:
                    output.append(f"\n... (terpotong, {ws.max_row-MAX_ROWS} baris lagi)")
                output.append("")

            if len(wb.sheetnames) > MAX_SHEETS:
                output.append(f"... (terpotong, {len(wb.sheetnames)-MAX_SHEETS} sheet lagi)")

            wb.close()
            return "\n".join(output)
        except ImportError:
            return "FAILED|openpyxl belum terinstall. Jalankan: pip install openpyxl"
        except Exception as e:
            return f"FAILED|Gagal baca Excel: {e}"

# ============================================================
# Tool 5: Baca Word (DOCX)
# ============================================================
class WordReader(BaseTool):
    name: str = "Word Reader Tool"
    description: str = """Baca dan ekstrak teks + tabel dari file Word (.docx).
    Input: path lokal file Word."""

    def _run(self, filepath: str) -> str:
        try:
            from docx import Document
            
            p = Path(filepath)
            if not p.exists():
                return f"FAILED|File tidak ditemukan: {filepath}"
            
            doc = Document(p)
            output = [f"=== Word: {p.name} ===\n"]
            
            # Metadata
            output.append(f"Total paragraf: {len(doc.paragraphs)}")
            output.append(f"Total tabel: {len(doc.tables)}\n")
            
            # Paragraf
            MAX_PARA = 100
            output.append("--- TEKS ---")
            for i, para in enumerate(doc.paragraphs[:MAX_PARA]):
                if para.text.strip():
                    output.append(para.text)

            if len(doc.paragraphs) > MAX_PARA:
                output.append(f"\n... (terpotong, {len(doc.paragraphs)-MAX_PARA} paragraf lagi)")

            # Tabel
            if doc.tables:
                output.append(f"\n--- TABEL ({len(doc.tables)}) ---")
                for i, table in enumerate(doc.tables[:5]):
                    output.append(f"\nTabel {i+1}:")
                    for row in table.rows[:15]:
                        cells = [cell.text for cell in row.cells]
                        output.append(f"  {' | '.join(cells)}")
                    if len(table.rows) > 15:
                        output.append(f"  ... ({len(table.rows)-15} baris lagi)")
                if len(doc.tables) > 5:
                    output.append(f"\n... ({len(doc.tables)-5} tabel lagi)")
            
            return "\n".join(output)
        except ImportError:
            return "FAILED|python-docx belum terinstall. Jalankan: pip install python-docx"
        except Exception as e:
            return f"FAILED|Gagal baca Word: {e}"

# ============================================================
# Tool: PowerPoint Reader (PPTX)
# ============================================================
class PowerPointReader(BaseTool):
    name: str = "PowerPoint Reader Tool"
    description: str = """Baca dan ekstrak teks + tabel dari file PowerPoint (.pptx).
    Input: path lokal file PowerPoint."""

    def _run(self, filepath: str) -> str:
        try:
            from pptx import Presentation

            p = Path(filepath)
            if not p.exists():
                return f"FAILED|File tidak ditemukan: {filepath}"

            prs = Presentation(p)
            output = [f"=== PowerPoint: {p.name} ===\n"]
            output.append(f"Total slide: {len(prs.slides)}\n")

            MAX_SLIDES = 30
            for i, slide in enumerate(prs.slides[:MAX_SLIDES]):
                output.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            txt = para.text.strip()
                            if txt:
                                output.append(txt)
                    if shape.has_table:
                        output.append("[TABEL]")
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            output.append("  " + " | ".join(cells))
                output.append("")

            if len(prs.slides) > MAX_SLIDES:
                output.append(f"... (terpotong, {len(prs.slides)-MAX_SLIDES} slide lagi)")

            return "\n".join(output)
        except ImportError:
            return "FAILED|python-pptx belum terinstall. Jalankan: pip install python-pptx"
        except Exception as e:
            return f"FAILED|Gagal baca PowerPoint: {e}"

# ============================================================
# Tool: Audio Transcriber (MP3/WAV/M4A/OGG)
# ============================================================
class AudioTranscriber(BaseTool):
    name: str = "Audio Transcriber Tool"
    description: str = """Transkripsikan file audio (MP3, WAV, M4A, OGG, OPUS, FLAC) menjadi teks.
    Cocok untuk voice note Discord, rekaman meeting, podcast.
    Output: transkrip lengkap + ringkasan poin penting.
    Input: path lokal file audio."""

    def _run(self, filepath: str) -> str:
        try:
            from openai import OpenAI

            p = Path(filepath)
            if not p.exists():
                return f"FAILED|File tidak ditemukan: {filepath}"

            file_size = p.stat().st_size
            if file_size > 25 * 1024 * 1024:
                return f"FAILED|Audio terlalu besar ({file_size // 1024 // 1024}MB). Maksimal 25MB."

            ext = p.suffix.lower().lstrip('.')
            if ext not in ('mp3', 'wav', 'm4a', 'ogg', 'opus', 'flac'):
                return f"FAILED|Format .{ext} tidak didukung. Pakai MP3/WAV/M4A/OGG/OPUS/FLAC."

            with open(p, 'rb') as f:
                audio_b64 = base64.b64encode(f.read()).decode('utf-8')

            client = OpenAI(
                api_key=os.environ.get("OPENROUTER_API_KEY"),
                base_url="https://openrouter.ai/api/v1"
            )

            result = client.chat.completions.create(
                model="google/gemini-2.5-flash",
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "input_audio",
                            "input_audio": {"data": audio_b64, "format": ext}
                        },
                        {
                            "type": "text",
                            "text": """Transkripsikan audio ini ke teks bahasa Indonesia.
Format output WAJIB:

=== TRANSKRIP ===
[transkrip lengkap kata per kata, pertahankan tanda baca alami]

=== RINGKASAN ===
- [3-5 poin penting dari isi audio]

=== METADATA ===
- Bahasa: [terdeteksi]
- Durasi estimasi: [pendek/sedang/panjang]
- Jumlah pembicara: [perkiraan]"""
                        }
                    ]
                }],
                max_tokens=4000
            )
            return result.choices[0].message.content
        except Exception as e:
            return f"FAILED|Gagal transkripsi audio: {e}"

# ============================================================
# Tool: Catalog Extractor (Furnitur Kompetitor → Structured JSON → Vault)
# ============================================================
class CatalogExtractorTool(BaseTool):
    name: str = "Catalog Extractor Tool"
    description: str = """Ekstrak data produk terstruktur dari katalog furnitur kompetitor (PDF brosur atau gambar).
    Output: tabel produk lengkap (nama, kode, harga, dimensi, material, kategori) dalam markdown.
    OTOMATIS arsip hasil ke vault Obsidian sebagai data referensi untuk riset kompetitor.
    Input: path lokal file katalog (PDF multi-halaman atau image PNG/JPG)."""

    def _run(self, filepath: str) -> str:
        try:
            from datetime import datetime as _dt
            import logging as _lg
            _logger = _lg.getLogger('bima_core')

            p = Path(filepath)
            if not p.exists():
                return f"FAILED|File tidak ditemukan: {filepath}"

            ext = p.suffix.lower()
            MAX_PAGES = 15
            pages_b64 = []

            # Build list of pages as base64 PNG
            if ext == '.pdf':
                import pdfplumber
                with pdfplumber.open(p) as pdf:
                    total_pages = len(pdf.pages)
                    for page in pdf.pages[:MAX_PAGES]:
                        img = page.to_image(resolution=150).original
                        buf = io.BytesIO()
                        img.save(buf, format='PNG')
                        pages_b64.append(base64.b64encode(buf.getvalue()).decode('utf-8'))
            elif ext in ('.png', '.jpg', '.jpeg', '.webp', '.bmp'):
                with open(p, 'rb') as f:
                    pages_b64.append(base64.b64encode(f.read()).decode('utf-8'))
                total_pages = 1
            else:
                return f"FAILED|Format {ext} tidak didukung. Pakai PDF atau gambar (PNG/JPG)."

            # Specialized prompt untuk katalog furnitur
            prompt = """Kamu adalah ahli analisis katalog furnitur kompetitor.
Ekstrak SEMUA produk furnitur yang terlihat dari halaman katalog ini.

Untuk SETIAP produk, ekstrak field berikut (kalau tidak ada di katalog → null):
- name: nama produk lengkap (string)
- code: kode/SKU produk (string atau null)
- price: harga + mata uang (string atau null, contoh: "Rp 2.499.000")
- dimensions: dimensi P×L×T (string atau null, contoh: "160x200x100 cm")
- material: material utama (string atau null, contoh: "MDF + veneer oak")
- color_options: opsi warna (array string, contoh: ["white","oak","black"])
- category: kategori (string: sofa/meja/kursi/tempat tidur/lemari/storage/lainnya)
- notes: diskon/bundle/garansi/promo (string atau null)

Output WAJIB JSON murni TANPA markdown code fence, struktur:
{
  "brand": "nama brand kalau terlihat di header/footer/logo, kalau tidak null",
  "page_summary": "1 kalimat ringkasan halaman ini",
  "products": [ { ...produk1... }, { ...produk2... } ]
}

ATURAN KETAT:
- Output HANYA JSON valid, tanpa penjelasan, tanpa ```json fence.
- Kalau halaman kosong/cover/index → products = [].
- Jangan ngarang produk yang tidak ada di katalog."""

            # Call Gemini Vision per halaman
            all_products = []
            brand = None
            page_summaries = []

            for i, b64 in enumerate(pages_b64):
                try:
                    raw = _gemini_vision(b64, "image/png", prompt, max_tokens=3000)
                    raw = raw.strip()
                    # Strip markdown fence kalau ada
                    if raw.startswith('```'):
                        raw = raw.split('\n', 1)[1] if '\n' in raw else raw[3:]
                        if raw.endswith('```'):
                            raw = raw.rsplit('```', 1)[0]
                    raw = raw.strip()
                    page_data = json.loads(raw)
                    if page_data.get('brand') and not brand:
                        brand = page_data['brand']
                    if page_data.get('page_summary'):
                        page_summaries.append(f"Hal {i+1}: {page_data['page_summary']}")
                    for prod in page_data.get('products', []) or []:
                        prod['source_page'] = i + 1
                        all_products.append(prod)
                except json.JSONDecodeError as e:
                    _logger.warning(f"[CATALOG] Parse JSON gagal di halaman {i+1}: {e}")
                except Exception as e:
                    _logger.warning(f"[CATALOG] Gemini call gagal di halaman {i+1}: {e}")

            if not all_products:
                return f"FAILED|Tidak ada produk berhasil diekstrak dari {p.name}. Pastikan file adalah katalog furnitur."

            # Format markdown untuk vault
            timestamp = _dt.now().strftime("%Y-%m-%d %H:%M")
            md_lines = [
                f"# Katalog Kompetitor: {brand or p.stem}",
                "",
                f"**Sumber file:** {p.name}",
                f"**Diekstrak:** {timestamp} WIB",
                f"**Jumlah produk:** {len(all_products)}",
                f"**Halaman diproses:** {len(pages_b64)} dari {total_pages}",
                "",
                "## Ringkasan Per Halaman",
                "",
                "\n".join(page_summaries) or "_(tidak ada ringkasan)_",
                "",
                "## Daftar Produk",
                "",
                "| # | Nama | Kode | Harga | Dimensi | Material | Kategori | Warna | Catatan | Hal |",
                "|---|------|------|-------|---------|----------|----------|-------|---------|-----|",
            ]
            for idx, prod in enumerate(all_products, 1):
                colors = prod.get('color_options') or []
                colors_str = ', '.join(colors) if isinstance(colors, list) else str(colors)
                row = [
                    str(idx),
                    str(prod.get('name', '-')),
                    str(prod.get('code') or '-'),
                    str(prod.get('price') or '-'),
                    str(prod.get('dimensions') or '-'),
                    str(prod.get('material') or '-'),
                    str(prod.get('category') or '-'),
                    colors_str or '-',
                    str(prod.get('notes') or '-'),
                    str(prod.get('source_page', '-')),
                ]
                # Sanitize untuk markdown table (hindari pipe & newline)
                row = [c.replace('|', '/').replace('\n', ' ') for c in row]
                md_lines.append("| " + " | ".join(row) + " |")

            md_lines.extend([
                "",
                "## Data JSON Mentah",
                "",
                "```json",
                json.dumps({"brand": brand, "products": all_products}, indent=2, ensure_ascii=False),
                "```",
                "",
                "*Diekstrak otomatis oleh CatalogExtractorTool — B.I.M.A Core*",
            ])
            markdown_content = "\n".join(md_lines)

            # Auto-archive ke vault via t3_arsip
            archive_status = ""
            try:
                from teams.t3_arsip import VaultSaveTool
                slug = (brand or p.stem).replace(' ', '_')[:40]
                vault_input = json.dumps({
                    "title": f"Katalog_{slug}_{_dt.now().strftime('%Y%m%d_%H%M')}",
                    "content": markdown_content
                })
                save_result = VaultSaveTool()._run(vault_input)
                archive_status = f"\n--- Auto-archive ---\n{save_result}"
            except Exception as e:
                archive_status = f"\n--- Auto-archive ---\nFAILED|{e}"

            # Return ringkasan ke agent (jangan paste seluruh markdown ke context LLM)
            preview_lines = [
                f"=== Catalog Extractor: {p.name} ===",
                f"Brand: {brand or '(tidak terdeteksi)'}",
                f"Total produk: {len(all_products)}",
                f"Halaman diproses: {len(pages_b64)} / {total_pages}",
                "",
                "--- Preview 10 Produk Pertama ---",
            ]
            for i, prod in enumerate(all_products[:10], 1):
                preview_lines.append(
                    f"{i}. {prod.get('name', '?')} | {prod.get('price', '-')} | "
                    f"{prod.get('dimensions', '-')} | hal {prod.get('source_page', '?')}"
                )
            if len(all_products) > 10:
                preview_lines.append(f"... dan {len(all_products)-10} produk lainnya (lihat di vault)")
            preview_lines.append(archive_status)

            return "\n".join(preview_lines)

        except Exception as e:
            return f"FAILED|Catalog Extractor error: {e}"

# ============================================================
# Tool 6: Analisis Gambar (sudah ada, diperkuat)
# ============================================================
class ImageAnalyzerTool(BaseTool):
    name: str = "Image Analyzer Tool"
    description: str = """Analisis gambar dari URL atau path lokal.
    Support: screenshot AutoCAD, D5 Render, foto material, katalog, dll.
    Input: URL gambar atau path lokal file gambar."""

    def _run(self, image_input: str) -> str:
        try:
            # Cek apakah input adalah path lokal atau URL
            p = Path(image_input)
            if p.exists():
                # Validasi ukuran file (max 10MB)
                file_size = p.stat().st_size
                if file_size > 10 * 1024 * 1024:
                    return f"FAILED|File terlalu besar ({file_size // 1024 // 1024}MB). Maksimal 10MB untuk analisis gambar."
                # File lokal
                with open(p, "rb") as f:
                    image_data = base64.b64encode(f.read()).decode("utf-8")
                suffix = p.suffix.lower()
                content_type = {
                    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                    ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp"
                }.get(suffix, "image/png")
            else:
                # URL
                resp = httpx.get(image_input, timeout=30)
                if len(resp.content) > 10 * 1024 * 1024:
                    return "FAILED|Gambar dari URL terlalu besar (>10MB)."
                image_data = base64.b64encode(resp.content).decode("utf-8")
                content_type = resp.headers.get("content-type", "image/png").split(";")[0]
            
            from openai import OpenAI
            client = OpenAI(
                api_key=os.environ.get("OPENROUTER_API_KEY"),
                base_url="https://openrouter.ai/api/v1"
            )
            
            result = client.chat.completions.create(
                model="google/gemini-3-flash-preview",
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{content_type};base64,{image_data}"
                            }
                        },
                        {
                            "type": "text",
                            "text": """Analisis gambar ini secara detail. Fokus pada:
                            1. Dimensi dan ukuran yang terlihat
                            2. Material dan spesifikasi teknis
                            3. Teks, label, atau angka yang ada
                            4. Elemen desain penting
                            Berikan hasil dalam Bahasa Indonesia."""
                        }
                    ]
                }],
                max_tokens=1500
            )
            return result.choices[0].message.content
        except Exception as e:
            return f"Gagal analisis gambar: {e}"

# ============================================================
# Tool 7: Image to Code (UI/UX Generator)
# ============================================================
class ImageToCodeTool(BaseTool):
    name: str = "Image to Code Tool"
    description: str = """Konversi gambar desain UI/UX (screenshot web, mockup) menjadi kode HTML & Tailwind CSS murni.
    Input: URL gambar atau path lokal file gambar mock up."""

    def _run(self, image_input: str) -> str:
        try:
            p = Path(image_input)
            if p.exists():
                file_size = p.stat().st_size
                if file_size > 10 * 1024 * 1024:
                    return f"FAILED|File terlalu besar ({file_size // 1024 // 1024}MB)."
                with open(p, "rb") as f:
                    image_data = base64.b64encode(f.read()).decode("utf-8")
                suffix = p.suffix.lower()
                content_type = {
                    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                    ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp"
                }.get(suffix, "image/png")
            else:
                resp = httpx.get(image_input, timeout=30)
                if len(resp.content) > 10 * 1024 * 1024:
                    return "FAILED|Gambar dari URL terlalu besar (>10MB)."
                image_data = base64.b64encode(resp.content).decode("utf-8")
                content_type = resp.headers.get("content-type", "image/png").split(";")[0]
            
            from openai import OpenAI
            client = OpenAI(
                api_key=os.environ.get("OPENROUTER_API_KEY"),
                base_url="https://openrouter.ai/api/v1"
            )
            
            result = client.chat.completions.create(
                model="google/gemini-3-flash-preview",
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{content_type};base64,{image_data}"
                            }
                        },
                        {
                            "type": "text",
                            "text": """Konversi gambar desain UI ini menjadi kode HTML5 dan Tailwind CSS yang rapi dan responsif.
                            ATURAN:
                            1. Gunakan semantic HTML5 (header, main, section, footer).
                            2. Gunakan CDN Tailwind CSS: <script src="https://cdn.tailwindcss.com"></script>
                            3. Buat semirip mungkin dengan gambar asli dari segi layout, warna, dan tipografi.
                            4. KEMBALIKAN HANYA KODE HTML PENUH (dimulai dengan <!DOCTYPE html> dan diakhiri </html>). Tanpa penjelasan apapun."""
                        }
                    ]
                }],
                max_tokens=4000
            )
            
            code = result.choices[0].message.content.strip()
            if code.startswith("```html"):
                code = code[7:-3].strip()
            elif code.startswith("```"):
                code = code[3:-3].strip()

            import time, hashlib
            slug = hashlib.md5(image_input.encode()).hexdigest()[:8]
            filename = f"ui_generated_{slug}_{int(time.time())}.html"
            filepath = OUTPUT_DIR / filename
            filepath.write_text(code, encoding="utf-8")

            return f"SUCCESS|{filepath}|File HTML berhasil dibuat dari gambar: {filename}"
        except Exception as e:
            return f"Gagal konversi gambar ke kode: {e}"

# ============================================================
# Visual Agent (Team 2 — Dirombak Total)
# ============================================================
visual_agent = Agent(
    role='Universal Document & Media Inspector',
    goal='Membaca, menganalisis, dan mengekstrak data dari SEMUA jenis file: PDF (termasuk PDF scan via OCR), Excel, Word, PowerPoint, TXT, CSV, gambar, dan audio dengan akurasi tinggi.',
    backstory="""Kamu adalah mata dewa dari B.I.M.A Core.
    Kamu bisa membaca SEMUA jenis file yang Bima kirim:
    - PDF (termasuk SCAN) → pakai PDFReader (otomatis OCR via Gemini Vision kalau text layer kosong)
    - Excel/XLSX → pakai ExcelReader (multiple sheets, max 50 baris/sheet)
    - Word/DOCX → pakai WordReader (teks + semua tabel)
    - PowerPoint/PPTX → pakai PowerPointReader (semua slide + tabel)
    - TXT/CSV/MD → pakai TextFileReader
    - Audio (MP3/WAV/M4A/OGG/OPUS/FLAC) → pakai AudioTranscriber (transkrip + ringkasan)
    - Gambar → pakai ImageAnalyzerTool (analisis visual)
    - Konversi UI ke Code → pakai ImageToCodeTool (jika diminta membuat web dari gambar)
    - KATALOG FURNITUR KOMPETITOR → pakai CatalogExtractorTool (ekstrak produk + harga + dimensi → otomatis arsip ke vault sebagai data referensi)

    Workflow saat Bima kirim file:
    1. Download file dari URL pakai FileDownloader
    2. Identifikasi tipe file dari ekstensi DAN konteksnya
    3. Kalau Bima bilang "ini katalog kompetitor", "rangkum katalog ini", "data referensi furnitur" → WAJIB pakai CatalogExtractorTool (bukan PDFReader biasa)
    4. Untuk file lain → panggil tool reader yang sesuai
    5. Ekstrak semua data penting
    6. Laporkan ke Bima dengan rapi

    Kamu TIDAK PERNAH mengarang data yang tidak ada di file.
    Selalu sebutkan sumber halaman/sheet/slide/file tempat data ditemukan.""",
    llm=visual_llm,
    tools=[
        FileDownloader(),
        TextFileReader(),
        PDFReader(),
        ExcelReader(),
        WordReader(),
        PowerPointReader(),
        AudioTranscriber(),
        CatalogExtractorTool(),
        ImageAnalyzerTool(),
        ImageToCodeTool()
    ],
    allow_delegation=True,
    verbose=True
)